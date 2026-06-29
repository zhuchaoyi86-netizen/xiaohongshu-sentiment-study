from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path("/Users/xinxinhuashe/Documents/1/python课/小组共享材料/第四部分_第1页视觉稿_可编辑试还原.pptx")


NAVY = RGBColor(22, 45, 78)
INK = RGBColor(54, 64, 77)
MUTED = RGBColor(118, 126, 136)
RED = RGBColor(235, 83, 73)
RED_SOFT = RGBColor(255, 231, 226)
ORANGE = RGBColor(238, 139, 66)
ORANGE_SOFT = RGBColor(255, 239, 221)
TEAL = RGBColor(42, 145, 139)
TEAL_SOFT = RGBColor(222, 244, 241)
BLUE = RGBColor(67, 118, 165)
BLUE_SOFT = RGBColor(229, 238, 249)
CREAM = RGBColor(253, 249, 243)
PAPER = RGBColor(255, 255, 252)
LINE = RGBColor(229, 221, 212)
WHITE = RGBColor(255, 255, 255)


def add_text(slide, x, y, w, h, text, size=14, bold=False, color=INK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "PingFang SC"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def rect(slide, x, y, w, h, fill, line=None, radius=True):
    kind = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shp = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.9)
    return shp


def oval(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.8)
    return shp


def connector(slide, x1, y1, x2, y2, color=LINE, width=1.2):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(width)
    return line


def dot_grid(slide, x, y, rows=5, cols=7, color=RGBColor(235, 228, 220)):
    for r in range(rows):
        for c in range(cols):
            oval(slide, x + c * 0.14, y + r * 0.14, 0.025, 0.025, color)


def icon_label(slide, x, y, icon, title, note, fill=WHITE, accent=BLUE):
    rect(slide, x, y, 1.1, 0.62, fill, LINE)
    oval(slide, x + 0.12, y + 0.16, 0.3, 0.3, accent)
    add_text(slide, x + 0.12, y + 0.16, 0.3, 0.3, icon, 9, True, WHITE, PP_ALIGN.CENTER)
    add_text(slide, x + 0.48, y + 0.12, 0.48, 0.16, title, 6.5, True, NAVY)
    add_text(slide, x + 0.48, y + 0.34, 0.48, 0.12, note, 5.2, False, MUTED)


def food_card(slide, x, y, w=2.28, h=3.1):
    card = rect(slide, x, y, w, h, WHITE, RGBColor(236, 226, 217))
    card.rotation = -4
    rect(slide, x + 0.18, y + 0.2, w - 0.36, 1.45, ORANGE_SOFT, None)
    # Plate and food marks
    oval(slide, x + 0.72, y + 0.38, 0.78, 0.58, WHITE, RGBColor(246, 196, 160))
    for dx, dy, col in [(0.86, 0.48, RED), (1.05, 0.6, ORANGE), (1.2, 0.48, TEAL), (1.02, 0.42, RGBColor(250, 190, 105))]:
        oval(slide, x + dx, y + dy, 0.12, 0.12, col)
    add_text(slide, x + 0.25, y + 1.82, w - 0.5, 0.18, "食品种草笔记", 9, True, NAVY)
    add_text(slide, x + 0.25, y + 2.12, w - 0.5, 0.28, "情绪密度 · 混合评价 · 点赞热度", 6.2, False, MUTED)
    oval(slide, x + 0.25, y + 2.67, 0.28, 0.28, RED_SOFT)
    add_text(slide, x + 0.25, y + 2.67, 0.28, 0.28, "♥", 10, True, RED, PP_ALIGN.CENTER)
    add_text(slide, x + 0.62, y + 2.7, 0.6, 0.16, "5.2k", 8, True, RED)


def metric_pill(slide, x, y, icon, label, value):
    rect(slide, x, y, 1.5, 0.42, WHITE, RGBColor(237, 229, 220))
    oval(slide, x + 0.12, y + 0.11, 0.2, 0.2, BLUE_SOFT)
    add_text(slide, x + 0.12, y + 0.1, 0.2, 0.2, icon, 6.5, True, BLUE, PP_ALIGN.CENTER)
    add_text(slide, x + 0.4, y + 0.11, 0.58, 0.12, label, 6.5, True, MUTED)
    add_text(slide, x + 0.94, y + 0.11, 0.38, 0.12, value, 7.2, True, NAVY)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background and right soft food atmosphere
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = CREAM
    rect(slide, 0, 0, 13.333, 7.5, RGBColor(255, 253, 248), None, radius=False)
    oval(slide, 10.95, 0.15, 2.5, 2.5, RGBColor(255, 245, 232))
    oval(slide, 11.65, 5.35, 1.75, 1.75, RGBColor(255, 239, 221))
    dot_grid(slide, 10.4, 0.78, rows=6, cols=7)

    # Left title block
    add_text(slide, 0.66, 0.62, 1.35, 0.22, "MODULE 04", 8.5, True, MUTED)
    rect(slide, 0.68, 0.92, 0.38, 0.035, RED, None, radius=False)
    add_text(slide, 0.66, 1.38, 3.65, 1.18, "模块四：\n实证模型与数据分析", 28, True, NAVY)
    rect(slide, 0.67, 3.18, 0.045, 0.74, RED, None, radius=False)
    add_text(slide, 0.86, 3.14, 3.35, 0.52, "基于文本挖掘的小红书食品种草笔记\n情感对点赞热度的影响研究", 10.5, True, RED)

    # Bottom information rail
    metric_pill(slide, 0.65, 6.28, "样", "样本量", "2899")
    metric_pill(slide, 2.35, 6.28, "OLS", "方法", "回归")
    metric_pill(slide, 4.05, 6.28, "✓", "检验", "鲁棒")
    metric_pill(slide, 5.75, 6.28, "图", "展示", "可视化")

    # Central model diagram
    rect(slide, 5.55, 1.05, 3.8, 4.5, WHITE, RGBColor(239, 229, 220))
    add_text(slide, 6.0, 1.42, 2.95, 0.22, "变量如何影响点赞热度？", 13.5, True, NAVY, PP_ALIGN.CENTER)
    nodes = [
        (6.0, 2.15, "情绪密度", TEAL_SOFT, TEAL),
        (6.0, 3.15, "混合评价", RED_SOFT, RED),
        (7.72, 2.65, "点赞热度", ORANGE_SOFT, RED),
    ]
    for x, y, label, fill, color in nodes:
        rect(slide, x, y, 1.18, 0.45, fill, fill)
        add_text(slide, x, y, 1.18, 0.45, label, 9.5, True, color, PP_ALIGN.CENTER)
    connector(slide, 7.18, 2.38, 7.72, 2.88, TEAL, 1.8)
    connector(slide, 7.18, 3.38, 7.72, 3.02, RED, 1.8)
    add_text(slide, 5.98, 4.12, 2.95, 0.45, "情绪表达越集中、评价态度越清晰，越容易获得点赞。", 8.2, False, INK, PP_ALIGN.CENTER)

    # Right social note card and small labels
    food_card(slide, 10.0, 1.2)
    icon_label(slide, 8.95, 1.58, "♥", "点赞热度", "log_likes", RED_SOFT, RED)
    icon_label(slide, 8.74, 2.72, "↗", "情绪密度", "核心变量", TEAL_SOFT, TEAL)
    icon_label(slide, 8.92, 3.85, "±", "混合评价", "负向作用", RED_SOFT, RED)

    # Small chart fragments near lower right
    rect(slide, 9.15, 5.2, 1.55, 0.88, WHITE, RGBColor(237, 229, 220))
    for i, h in enumerate([0.28, 0.42, 0.62, 0.36]):
        rect(slide, 9.38 + i * 0.27, 5.85 - h, 0.13, h, [BLUE_SOFT, TEAL_SOFT, ORANGE_SOFT, RED_SOFT][i], None, radius=False)
    connector(slide, 5.8, 5.35, 9.0, 5.42, RED_SOFT, 1.2)
    connector(slide, 5.4, 4.75, 4.68, 5.56, RED_SOFT, 1.0)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
