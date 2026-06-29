from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path("/Users/xinxinhuashe/Documents/1/python课")
PART4_DIR = ROOT / "第四部分材料"
FIG_DIR = PART4_DIR / "图"
TAB_DIR = PART4_DIR / "表"
OUT_PPT = ROOT / "小组共享材料" / "第四部分_实证模型与数据分析.pptx"

BG = RGBColor(247, 243, 236)
PANEL = RGBColor(255, 252, 247)
NAVY = RGBColor(23, 49, 77)
BLUE = RGBColor(76, 122, 164)
BLUE_LIGHT = RGBColor(213, 229, 242)
ORANGE = RGBColor(196, 124, 56)
ORANGE_LIGHT = RGBColor(244, 226, 205)
RED = RGBColor(168, 73, 73)
GREEN = RGBColor(86, 127, 92)
TEXT = RGBColor(51, 58, 68)
MUTED = RGBColor(118, 125, 132)
WHITE = RGBColor(255, 255, 255)


def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_text(slide, left, top, width, height, text, size=18, bold=False, color=TEXT, align=PP_ALIGN.LEFT, font="PingFang SC"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_multiline(slide, left, top, width, height, lines, size=12, color=TEXT, bullet=True, spacing=5):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {line}" if bullet else line
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(spacing)
        for run in p.runs:
            run.font.name = "PingFang SC"
            run.font.size = Pt(size)
            run.font.color.rgb = color
    return box


def card(slide, left, top, width, height, fill=PANEL, line=RGBColor(228, 220, 210), radius=True):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(1.0)
    return shp


def chip(slide, left, top, width, height, text, fill=BLUE_LIGHT, color=NAVY):
    shp = card(slide, left, top, width, height, fill=fill, line=fill)
    add_text(slide, left, top + height * 0.18, width, height * 0.64, text, size=11, bold=True, color=color, align=PP_ALIGN.CENTER)
    return shp


def section_header(slide, title, subtitle=None, number=None):
    set_bg(slide)
    card(slide, Inches(0.55), Inches(0.35), Inches(12.2), Inches(0.95), fill=PANEL)
    if number:
        tag = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.56), Inches(0.7), Inches(0.38))
        tag.fill.solid()
        tag.fill.fore_color.rgb = ORANGE
        tag.line.fill.background()
        add_text(slide, Inches(0.8), Inches(0.6), Inches(0.7), Inches(0.2), number, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        x = Inches(1.7)
    else:
        x = Inches(0.8)
    add_text(slide, x, Inches(0.56), Inches(8.8), Inches(0.32), title, size=24, bold=True, color=NAVY)
    if subtitle:
        add_text(slide, x, Inches(0.93), Inches(9.8), Inches(0.2), subtitle, size=10, color=MUTED)


def add_picture(slide, path, left, top, width=None, height=None):
    kwargs = {}
    if width is not None:
        kwargs["width"] = width
    if height is not None:
        kwargs["height"] = height
    slide.shapes.add_picture(str(path), left, top, **kwargs)


def metric_card(slide, left, top, title, value, note, tone="blue"):
    fill = BLUE_LIGHT if tone == "blue" else ORANGE_LIGHT
    accent = BLUE if tone == "blue" else ORANGE
    card(slide, left, top, Inches(2.45), Inches(1.25), fill=fill, line=fill)
    add_text(slide, left + Inches(0.18), top + Inches(0.16), Inches(2.0), Inches(0.2), title, size=10, bold=True, color=accent)
    add_text(slide, left + Inches(0.18), top + Inches(0.42), Inches(2.0), Inches(0.36), value, size=22, bold=True, color=NAVY)
    add_text(slide, left + Inches(0.18), top + Inches(0.87), Inches(2.05), Inches(0.18), note, size=9, color=MUTED)


def add_table_simple(slide, df, left, top, width, height, font_size=10):
    rows, cols = df.shape[0] + 1, df.shape[1]
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    for j, col in enumerate(df.columns):
        cell = table.cell(0, j)
        cell.text = str(col)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.name = "PingFang SC"
                r.font.size = Pt(font_size)
                r.font.bold = True
                r.font.color.rgb = WHITE
    for i in range(df.shape[0]):
        for j in range(cols):
            cell = table.cell(i + 1, j)
            cell.text = str(df.iloc[i, j])
            cell.fill.solid()
            cell.fill.fore_color.rgb = PANEL if i % 2 == 0 else BG
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for r in p.runs:
                    r.font.name = "PingFang SC"
                    r.font.size = Pt(font_size)
                    r.font.color.rgb = TEXT
    return table


def connector(slide, x1, y1, x2, y2, color=MUTED):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(1.7)
    return line


def build_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    key_df = pd.read_csv(TAB_DIR / "第四部分_基准回归关键结果.csv")
    robust_density = pd.read_csv(TAB_DIR / "第四部分_鲁棒性_情绪密度.csv")
    robust_mixed = pd.read_csv(TAB_DIR / "第四部分_鲁棒性_混合评价.csv")
    summary_text = (PART4_DIR / "第四部分_结论摘要.txt").read_text(encoding="utf-8").splitlines()

    core_df = key_df[key_df["变量"].isin(["情绪密度", "混合评价虚拟变量", "配图数量", "感叹号数量", "上午发布", "下午发布", "晚上发布"])].copy()

    # 1 cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.0), Inches(0.0), Inches(4.2), Inches(7.5))
    accent.fill.solid()
    accent.fill.fore_color.rgb = NAVY
    accent.line.fill.background()
    add_text(slide, Inches(0.65), Inches(0.9), Inches(2.9), Inches(0.38), "模块四", size=16, bold=True, color=ORANGE_LIGHT)
    add_text(slide, Inches(0.65), Inches(1.45), Inches(3.0), Inches(1.9), "实证模型\n与数据分析", size=28, bold=True, color=WHITE)
    add_text(slide, Inches(0.65), Inches(3.75), Inches(2.8), Inches(0.85), "把“变量关系、鲁棒性检验、图形化展示”整合成一套可汇报的研究证据。", size=12, color=RGBColor(228, 234, 239))
    card(slide, Inches(4.55), Inches(0.7), Inches(8.0), Inches(5.95), fill=PANEL)
    add_text(slide, Inches(5.0), Inches(1.0), Inches(7.0), Inches(0.55), "基于文本挖掘的小红书食品种草笔记情感对点赞热度的影响研究", size=22, bold=True, color=NAVY)
    add_text(slide, Inches(5.0), Inches(1.72), Inches(5.8), Inches(0.32), "第四部分展示逻辑", size=11, bold=True, color=BLUE)
    chip(slide, Inches(5.0), Inches(2.1), Inches(1.45), Inches(0.42), "模型设定", fill=BLUE_LIGHT)
    chip(slide, Inches(6.62), Inches(2.1), Inches(1.62), Inches(0.42), "基准回归", fill=ORANGE_LIGHT, color=ORANGE)
    chip(slide, Inches(8.45), Inches(2.1), Inches(1.78), Inches(0.42), "变量关系", fill=BLUE_LIGHT)
    chip(slide, Inches(10.46), Inches(2.1), Inches(1.85), Inches(0.42), "鲁棒性检验", fill=ORANGE_LIGHT, color=ORANGE)
    add_multiline(slide, Inches(5.05), Inches(2.9), Inches(6.8), Inches(1.4), [
        "核心问题不是“帖子越正向就越容易火”，而是情绪表达方式、评价结构和发布时间如何共同影响点赞热度。",
        "因此第四部分重点突出两个更有研究味道的结论：情绪密度的正向作用、混合评价的稳定抑制作用。"
    ], size=13, bullet=False, spacing=8)
    metric_card(slide, Inches(5.0), Inches(4.6), "样本量", "2899", "食品种草笔记有效样本", tone="blue")
    metric_card(slide, Inches(7.65), Inches(4.6), "调整 R²", "0.1596", "基准模型解释力", tone="orange")
    metric_card(slide, Inches(10.3), Inches(4.6), "核心发现", "2 个", "情绪密度、混合评价", tone="blue")

    # 2 model setup
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    section_header(slide, "4.1 研究模型设定与变量关系", "先说明变量如何进入模型，再解释变量之间的作用逻辑。", "01")
    card(slide, Inches(0.7), Inches(1.45), Inches(3.4), Inches(5.4))
    add_text(slide, Inches(0.95), Inches(1.72), Inches(2.7), Inches(0.3), "变量设定", size=15, bold=True, color=NAVY)
    add_multiline(slide, Inches(0.95), Inches(2.15), Inches(2.9), Inches(4.2), [
        "因变量：缩尾后的对数点赞量 `log_likes_winsorized`。",
        "核心变量：情绪密度、情感得分、混合评价。",
        "辅助变量：情感方向虚拟变量、文本长度、配图数量、话题标签、标点语气。",
        "时间变量：以“凌晨”为基准组，加入上午、下午、晚上虚拟变量。"
    ], size=12)
    # framework
    for left, top, text, fill in [
        (Inches(4.65), Inches(1.8), "情绪密度", BLUE_LIGHT),
        (Inches(4.65), Inches(3.05), "情感得分", BLUE_LIGHT),
        (Inches(4.65), Inches(4.3), "混合评价", ORANGE_LIGHT),
        (Inches(8.9), Inches(3.05), "点赞热度", RGBColor(230, 238, 245)),
    ]:
        card(slide, left, top, Inches(2.25), Inches(0.85), fill=fill, line=fill)
        add_text(slide, left, top + Inches(0.23), Inches(2.25), Inches(0.2), text, size=15, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    connector(slide, Inches(6.9), Inches(2.22), Inches(8.9), Inches(3.45))
    connector(slide, Inches(6.9), Inches(3.47), Inches(8.9), Inches(3.47))
    connector(slide, Inches(6.9), Inches(4.72), Inches(8.9), Inches(3.47))
    card(slide, Inches(7.25), Inches(5.0), Inches(4.7), Inches(1.45), fill=RGBColor(243, 247, 249))
    add_text(slide, Inches(7.5), Inches(5.22), Inches(4.2), Inches(0.22), "变量关系要点", size=13, bold=True, color=BLUE)
    add_multiline(slide, Inches(7.48), Inches(5.55), Inches(4.1), Inches(0.7), [
        "情绪密度预计正向影响点赞热度。",
        "混合评价预计负向影响点赞热度。",
        "图文呈现和发布时间用于控制非文本干扰。"
    ], size=11)

    # 3 baseline results
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    section_header(slide, "4.2 基准线性回归模型结果", "不在台上堆大表，直接抓住最该讲的结果。", "02")
    metric_card(slide, Inches(0.8), Inches(1.45), "情绪密度系数", "0.0949", "方向为正，普通标准误下显著", tone="blue")
    metric_card(slide, Inches(3.55), Inches(1.45), "混合评价系数", "-0.3389", "显著为负，抑制点赞热度", tone="orange")
    metric_card(slide, Inches(6.3), Inches(1.45), "配图数量系数", "0.2386", "图文丰富度影响明显", tone="blue")
    metric_card(slide, Inches(9.05), Inches(1.45), "营销话术系数", "1.5046", "营销型表达显著抬升热度", tone="orange")
    add_picture(slide, FIG_DIR / "第四部分_核心变量系数图.png", Inches(0.85), Inches(3.0), width=Inches(5.8))
    card(slide, Inches(6.95), Inches(2.95), Inches(5.4), Inches(3.1))
    add_text(slide, Inches(7.2), Inches(3.2), Inches(4.8), Inches(0.25), "怎么解释这组结果", size=14, bold=True, color=NAVY)
    add_multiline(slide, Inches(7.2), Inches(3.6), Inches(4.7), Inches(2.1), [
        "情绪密度为正，说明情绪表达越集中、越鲜明，帖子越容易获得点赞。",
        "混合评价显著为负，说明一篇帖子如果同时传递推荐与保留意见，用户互动会更谨慎。",
        "配图数量和发布时间同样重要，说明小红书热度不是由文本单独决定的。"
    ], size=12)

    # 4 variable relation visual
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    section_header(slide, "4.3 变量之间的模型关系", "让老师一眼看到“变量是怎么连起来的”。", "03")
    card(slide, Inches(0.7), Inches(1.45), Inches(5.95), Inches(4.85))
    add_text(slide, Inches(0.95), Inches(1.7), Inches(3.2), Inches(0.22), "情绪密度与点赞热度", size=14, bold=True, color=NAVY)
    add_picture(slide, FIG_DIR / "第四部分_散点图_情绪密度与点赞.png", Inches(0.95), Inches(2.05), width=Inches(5.45))
    card(slide, Inches(6.85), Inches(1.45), Inches(5.8), Inches(2.3))
    add_text(slide, Inches(7.1), Inches(1.72), Inches(3.4), Inches(0.22), "分组均值对比", size=14, bold=True, color=NAVY)
    add_picture(slide, FIG_DIR / "第四部分_柱状图_情绪密度与点赞.png", Inches(7.1), Inches(2.02), width=Inches(5.25))
    card(slide, Inches(6.85), Inches(4.0), Inches(5.8), Inches(2.3))
    add_text(slide, Inches(7.1), Inches(4.24), Inches(3.2), Inches(0.2), "关系总结", size=14, bold=True, color=NAVY)
    add_multiline(slide, Inches(7.12), Inches(4.62), Inches(5.05), Inches(1.2), [
        "散点图显示总体斜率向上，但离散较大，说明关系存在但不是单一决定因素。",
        "分组柱状图更直观地表明：高情绪密度组的平均点赞明显高于低情绪密度组。"
    ], size=12)

    # 5 robustness
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    section_header(slide, "4.4 鲁棒性检验", "不仅要有结论，还要证明这个结论经得住换模型和换样本。", "04")
    card(slide, Inches(0.75), Inches(1.45), Inches(3.35), Inches(4.95))
    add_text(slide, Inches(1.0), Inches(1.72), Inches(2.8), Inches(0.2), "检验路径", size=14, bold=True, color=NAVY)
    add_multiline(slide, Inches(1.0), Inches(2.08), Inches(2.75), Inches(3.9), [
        "换标准误：HC3 稳健标准误",
        "换因变量：点赞、评论、收藏",
        "换解释变量：原始情感得分、正负占比",
        "换样本：剔除极端值、凌晨样本、营销样本",
        "子样本复测：仅保留正向情绪样本"
    ], size=12)
    add_picture(slide, ROOT / "图" / "forest_mixed.png", Inches(4.35), Inches(1.55), width=Inches(4.0))
    add_picture(slide, ROOT / "图" / "forest_density.png", Inches(8.55), Inches(1.55), width=Inches(4.0))
    chip(slide, Inches(4.5), Inches(5.85), Inches(3.5), Inches(0.5), "混合评价：方向稳定为负", fill=ORANGE_LIGHT, color=RED)
    chip(slide, Inches(8.72), Inches(5.85), Inches(3.3), Inches(0.5), "情绪密度：方向稳定为正", fill=BLUE_LIGHT, color=BLUE)

    # 6 visual + time
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    section_header(slide, "4.5 图形化结果展示", "把变量关系落到更直观的用户行为节奏上。", "05")
    card(slide, Inches(0.78), Inches(1.45), Inches(7.25), Inches(4.95))
    add_picture(slide, FIG_DIR / "第四部分_折线图_发布时间与点赞.png", Inches(1.0), Inches(1.78), width=Inches(6.8))
    card(slide, Inches(8.35), Inches(1.45), Inches(4.15), Inches(4.95), fill=PANEL)
    add_text(slide, Inches(8.65), Inches(1.72), Inches(3.3), Inches(0.2), "画图得到的补充发现", size=14, bold=True, color=NAVY)
    add_multiline(slide, Inches(8.65), Inches(2.12), Inches(3.2), Inches(3.2), [
        "发布时间存在明显波动，下午和晚间更容易形成高点赞热度。",
        "这说明文本情感并不是唯一因素，平台活跃时段会放大内容传播效果。",
        "因此第四部分要把“情感作用”理解为嵌入发布时间和图文表现之中的综合机制。 "
    ], size=12)
    chip(slide, Inches(8.66), Inches(5.55), Inches(3.2), Inches(0.48), "文本变量 + 时间变量共同作用", fill=RGBColor(231, 239, 233), color=GREEN)

    # 7 final conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    section_header(slide, "4.6 第四部分结论", "最后一页只保留最值得你在台上强调的四点。", "06")
    card(slide, Inches(0.9), Inches(1.55), Inches(11.55), Inches(4.95), fill=PANEL)
    conclusions = [line.split(". ", 1)[1] if ". " in line else line for line in summary_text[1:5]]
    y = 1.95
    colors = [BLUE_LIGHT, ORANGE_LIGHT, BLUE_LIGHT, ORANGE_LIGHT]
    for i, c in enumerate(conclusions, 1):
        card(slide, Inches(1.2), Inches(y), Inches(10.95), Inches(0.75), fill=colors[i - 1], line=colors[i - 1])
        add_text(slide, Inches(1.45), Inches(y + 0.16), Inches(0.4), Inches(0.2), f"{i}", size=18, bold=True, color=NAVY)
        add_text(slide, Inches(1.9), Inches(y + 0.12), Inches(9.8), Inches(0.42), c, size=13, color=TEXT)
        y += 0.94
    chip(slide, Inches(4.15), Inches(6.75), Inches(5.1), Inches(0.42), "最适合汇报的主结论：混合评价更稳，情绪密度更有解释力", fill=RGBColor(238, 230, 218), color=NAVY)

    # 8 appendix-like evidence slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    section_header(slide, "附：基准模型关键结果速览", "这一页留作答疑备用，不建议全讲。", "A")
    small = core_df[["变量", "系数", "p值", "结论"]].copy()
    add_table_simple(slide, small, Inches(0.85), Inches(1.55), Inches(5.0), Inches(4.9), font_size=10)
    robust_small = robust_mixed[["检验设定", "系数", "p值"]].head(8).copy()
    add_table_simple(slide, robust_small, Inches(6.15), Inches(1.55), Inches(6.15), Inches(4.9), font_size=8)
    add_text(slide, Inches(0.92), Inches(6.65), Inches(11.0), Inches(0.3), "答疑时可以重点指出：混合评价在多数鲁棒性设定下都显著为负，这也是本研究最“稳”的经验发现。", size=11, color=MUTED)

    prs.save(OUT_PPT)
    print(f"已生成优化版 PPT：{OUT_PPT}")


if __name__ == "__main__":
    build_ppt()
