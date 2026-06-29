from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path("/Users/xinxinhuashe/Documents/1/python课")
OUT = ROOT / "小组共享材料" / "第四部分_实证模型与数据分析_汇报优化版.pptx"
FIG = ROOT / "第四部分材料" / "图"
OLD_FIG = ROOT / "图"
TAB = ROOT / "第四部分材料" / "表"


BG = RGBColor(252, 249, 243)
PAPER = RGBColor(255, 255, 252)
NAVY = RGBColor(22, 45, 78)
INK = RGBColor(45, 55, 68)
MUTED = RGBColor(111, 120, 130)
LINE = RGBColor(228, 220, 209)
RED = RGBColor(232, 83, 74)
RED_SOFT = RGBColor(255, 230, 224)
TEAL = RGBColor(45, 144, 139)
TEAL_SOFT = RGBColor(222, 244, 241)
ORANGE = RGBColor(238, 139, 66)
ORANGE_SOFT = RGBColor(255, 239, 221)
BLUE_SOFT = RGBColor(226, 236, 249)
GREEN_SOFT = RGBColor(230, 242, 225)
WHITE = RGBColor(255, 255, 255)


def bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG


def text(slide, x, y, w, h, s, size=16, bold=False, color=INK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = s
    run.font.name = "PingFang SC"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def bullets(slide, x, y, w, h, items, size=11, color=INK, dot=True):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ("• " if dot else "") + item
        p.space_after = Pt(5)
        for r in p.runs:
            r.font.name = "PingFang SC"
            r.font.size = Pt(size)
            r.font.color.rgb = color
    return box


def rect(slide, x, y, w, h, fill=PAPER, line=LINE, radius=True):
    kind = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shp = slide.shapes.add_shape(kind, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(1)
    return shp


def tag(slide, x, y, s, fill=RED, color=WHITE):
    rect(slide, x, y, Inches(0.42), Inches(0.32), fill=fill, line=fill)
    text(slide, x, y, Inches(0.42), Inches(0.32), s, size=11, bold=True, color=color, align=PP_ALIGN.CENTER)


def header(slide, num, title, subtitle=None):
    bg(slide)
    tag(slide, Inches(0.55), Inches(0.38), num)
    text(slide, Inches(1.05), Inches(0.33), Inches(7.8), Inches(0.42), title, size=23, bold=True, color=NAVY)
    if subtitle:
        text(slide, Inches(1.05), Inches(0.78), Inches(9.6), Inches(0.26), subtitle, size=10, color=MUTED)
    # small dotted ornament
    for i in range(6):
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(11.6 + 0.14 * (i % 3)), Inches(0.43 + 0.14 * (i // 3)), Inches(0.035), Inches(0.035))
        dot.fill.solid()
        dot.fill.fore_color.rgb = RED if i == 0 else LINE
        dot.line.fill.background()


def metric(slide, x, y, label, value, note, fill=BLUE_SOFT, accent=TEAL):
    rect(slide, x, y, Inches(2.65), Inches(1.15), fill=fill, line=fill)
    text(slide, x + Inches(0.18), y + Inches(0.13), Inches(2.2), Inches(0.18), label, size=9, bold=True, color=accent)
    text(slide, x + Inches(0.18), y + Inches(0.39), Inches(2.1), Inches(0.34), value, size=21, bold=True, color=NAVY)
    text(slide, x + Inches(0.18), y + Inches(0.78), Inches(2.25), Inches(0.18), note, size=8, color=MUTED)


def connector(slide, x1, y1, x2, y2, color=MUTED, width=1.5):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    c.line.color.rgb = color
    c.line.width = Pt(width)
    return c


def picture(slide, path, x, y, w=None, h=None):
    kwargs = {}
    if w:
        kwargs["width"] = w
    if h:
        kwargs["height"] = h
    slide.shapes.add_picture(str(path), x, y, **kwargs)


def icon_circle(slide, x, y, label, fill, color=NAVY):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x, y, Inches(0.42), Inches(0.42))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    text(slide, x, y, Inches(0.42), Inches(0.42), label, size=12, bold=True, color=color, align=PP_ALIGN.CENTER)


def food_card(slide, x, y, w=2.3, h=2.75):
    rect(slide, x, y, Inches(w), Inches(h), fill=WHITE, line=RGBColor(235, 226, 216))
    pic = rect(slide, x + Inches(0.18), y + Inches(0.18), Inches(w - 0.36), Inches(1.2), fill=ORANGE_SOFT, line=ORANGE_SOFT)
    # simple plate
    plate = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x + Inches(0.72), y + Inches(0.34), Inches(0.82), Inches(0.58))
    plate.fill.solid()
    plate.fill.fore_color.rgb = WHITE
    plate.line.color.rgb = RGBColor(246, 198, 160)
    for dx, dy, col in [(0.82, 0.44, RED), (1.02, 0.56, ORANGE), (1.2, 0.43, TEAL)]:
        d = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x + Inches(dx), y + Inches(dy), Inches(0.12), Inches(0.12))
        d.fill.solid()
        d.fill.fore_color.rgb = col
        d.line.fill.background()
    text(slide, x + Inches(0.2), y + Inches(1.55), Inches(w - 0.4), Inches(0.22), "食品种草笔记", size=10, bold=True, color=NAVY)
    text(slide, x + Inches(0.2), y + Inches(1.86), Inches(w - 0.4), Inches(0.34), "情绪密度 · 点赞热度", size=8, color=MUTED)
    icon_circle(slide, x + Inches(0.22), y + Inches(2.32), "♥", RED_SOFT, RED)
    text(slide, x + Inches(0.72), y + Inches(2.32), Inches(1.0), Inches(0.22), "5.2k", size=10, bold=True, color=RED)


def slide_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    left = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(4.1), Inches(7.5))
    left.fill.solid()
    left.fill.fore_color.rgb = NAVY
    left.line.fill.background()
    text(s, Inches(0.65), Inches(0.75), Inches(1.3), Inches(0.28), "MODULE 04", size=11, bold=True, color=ORANGE_SOFT)
    text(s, Inches(0.65), Inches(1.38), Inches(3.0), Inches(1.65), "模块四：\n实证模型与数据分析", size=26, bold=True, color=WHITE)
    text(s, Inches(0.68), Inches(3.3), Inches(2.75), Inches(0.8), "基于文本挖掘的小红书食品种草笔记情感对点赞热度的影响研究", size=11, color=RGBColor(235, 240, 245))
    metric(s, Inches(0.65), Inches(5.45), "样本量", "2899", "食品种草笔记", fill=RGBColor(30, 61, 96), accent=ORANGE_SOFT)
    metric(s, Inches(2.28), Inches(5.45), "方法", "OLS", "稳健性检验", fill=RGBColor(30, 61, 96), accent=ORANGE_SOFT)
    rect(s, Inches(4.65), Inches(0.7), Inches(7.9), Inches(5.95), fill=PAPER)
    food_card(s, Inches(9.45), Inches(1.1), 2.45, 2.8)
    # model mini map
    nodes = [(5.25, 2.0, "情绪密度", TEAL_SOFT), (5.25, 3.25, "混合评价", RED_SOFT), (7.15, 2.65, "点赞热度", ORANGE_SOFT)]
    for x, y, lab, fill in nodes:
        rect(s, Inches(x), Inches(y), Inches(1.35), Inches(0.5), fill=fill, line=fill)
        text(s, Inches(x), Inches(y), Inches(1.35), Inches(0.5), lab, size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    connector(s, Inches(6.6), Inches(2.25), Inches(7.15), Inches(2.9), TEAL, 2)
    connector(s, Inches(6.6), Inches(3.5), Inches(7.15), Inches(3.05), RED, 2)
    text(s, Inches(5.0), Inches(4.6), Inches(4.0), Inches(0.35), "情绪密度 / 混合评价 / 情感方向", size=18, bold=True, color=NAVY)
    text(s, Inches(5.0), Inches(5.05), Inches(4.1), Inches(0.5), "通过文本表达影响点赞热度，同时受到配图、文本长度、发布时间等变量控制。", size=11, color=INK)
    return s


def slide_framework(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "4.1", "实证分析思路", "变量关系 → 模型检验 → 鲁棒性 → 可视化")
    text(s, Inches(0.85), Inches(1.35), Inches(6.8), Inches(0.45), "第四部分不是“堆模型”，而是建立一条从变量到结论的证据链。", size=16, bold=True, color=NAVY)
    items = [
        ("变量设定", "明确因变量、自变量、控制变量"),
        ("基准回归", "估计核心情感变量的影响"),
        ("变量关系", "解释情绪、评价结构和热度的联系"),
        ("鲁棒性检验", "换变量、换样本、换模型复测"),
        ("可视化展示", "用图形呈现主要关系"),
    ]
    start_x = 0.75
    gap = 2.45
    for i, (title, note) in enumerate(items):
        x = Inches(start_x + i * gap)
        rect(s, x, Inches(2.35), Inches(1.85), Inches(2.45), fill=WHITE)
        icon_circle(s, x + Inches(0.72), Inches(2.58), f"{i+1}", RED_SOFT if i in [1, 3] else BLUE_SOFT, RED if i in [1, 3] else NAVY)
        text(s, x + Inches(0.15), Inches(3.25), Inches(1.55), Inches(0.22), title, size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        text(s, x + Inches(0.18), Inches(3.68), Inches(1.5), Inches(0.52), note, size=9, color=MUTED, align=PP_ALIGN.CENTER)
        if i < 4:
            connector(s, x + Inches(1.85), Inches(3.58), x + Inches(2.43), Inches(3.58), ORANGE, 1.8)
    metric(s, Inches(0.85), Inches(5.65), "研究样本", "2899", "有效食品种草笔记", fill=ORANGE_SOFT, accent=RED)
    food_card(s, Inches(10.35), Inches(4.85), 1.65, 1.95)
    return s


def slide_model(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "4.2", "研究模型设定", "以点赞热度为因变量，以情感变量为核心自变量。")
    groups = [
        ("因变量", "log_likes_winsorized\n缩尾后的对数点赞量", RED_SOFT, RED),
        ("核心自变量", "sentiment_density\nsentiment_score_winsorized\nmixed_review_flag", BLUE_SOFT, NAVY),
        ("控制变量", "text_length_winsorized\nimagenumber / hashtag_count\ntime_period / marketing_flag", GREEN_SOFT, GREEN_SOFT),
    ]
    for i, (title, body, fill, accent) in enumerate(groups):
        y = Inches(1.45 + i * 1.5)
        rect(s, Inches(0.8), y, Inches(4.25), Inches(1.18), fill=fill, line=fill)
        text(s, Inches(1.05), y + Inches(0.14), Inches(1.1), Inches(0.22), title, size=12, bold=True, color=RED if i == 0 else NAVY)
        text(s, Inches(1.05), y + Inches(0.48), Inches(3.6), Inches(0.45), body, size=10, color=INK)
    rect(s, Inches(5.55), Inches(1.55), Inches(6.85), Inches(2.15), fill=WHITE)
    text(s, Inches(5.9), Inches(1.88), Inches(5.9), Inches(0.3), "基准模型", size=16, bold=True, color=NAVY)
    formula = "log_likes_winsorized = β0 + β1 情绪密度\n+ β2 情感得分 + β3 混合评价 + β4 Controls + ε"
    text(s, Inches(5.9), Inches(2.38), Inches(5.8), Inches(0.9), formula, size=16, bold=True, color=INK)
    rect(s, Inches(5.55), Inches(4.15), Inches(6.85), Inches(1.55), fill=ORANGE_SOFT, line=ORANGE_SOFT)
    text(s, Inches(5.9), Inches(4.4), Inches(5.9), Inches(0.28), "变量处理说明", size=14, bold=True, color=RED)
    text(s, Inches(5.9), Inches(4.82), Inches(5.8), Inches(0.45), "点赞量先缩尾再取对数，用来减弱爆款笔记极端值对估计结果的影响。", size=11, color=INK)
    return s


def slide_relation(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "4.3", "变量之间的模型关系", "点赞热度由情绪表达、评价结构和内容呈现共同影响。")
    center = (Inches(7.3), Inches(3.0))
    rect(s, center[0] - Inches(1.0), center[1] - Inches(0.55), Inches(2.05), Inches(1.1), fill=ORANGE_SOFT, line=ORANGE_SOFT)
    text(s, center[0] - Inches(0.85), center[1] - Inches(0.25), Inches(1.7), Inches(0.28), "点赞热度", size=18, bold=True, color=RED, align=PP_ALIGN.CENTER)
    text(s, center[0] - Inches(0.85), center[1] + Inches(0.12), Inches(1.7), Inches(0.2), "log_likes", size=9, color=MUTED, align=PP_ALIGN.CENTER)
    nodes = [
        (1.0, 1.65, "情绪密度", "正向促进", TEAL_SOFT, TEAL),
        (1.0, 3.0, "情感得分", "方向需结合密度解释", BLUE_SOFT, NAVY),
        (1.0, 4.35, "混合评价", "负向抑制", RED_SOFT, RED),
        (5.1, 5.6, "控制变量", "配图、文本、时间、营销话术", WHITE, MUTED),
    ]
    for x, y, title, note, fill, color in nodes:
        rect(s, Inches(x), Inches(y), Inches(2.55), Inches(0.85), fill=fill, line=fill if fill != WHITE else LINE)
        text(s, Inches(x + 0.18), Inches(y + 0.12), Inches(2.1), Inches(0.2), title, size=13, bold=True, color=color)
        text(s, Inches(x + 0.18), Inches(y + 0.43), Inches(2.1), Inches(0.2), note, size=9, color=MUTED)
    connector(s, Inches(3.55), Inches(2.08), Inches(6.3), Inches(3.15), TEAL, 2)
    connector(s, Inches(3.55), Inches(3.42), Inches(6.3), Inches(3.1), NAVY, 1.6)
    connector(s, Inches(3.55), Inches(4.77), Inches(6.3), Inches(3.2), RED, 2)
    connector(s, Inches(6.38), Inches(5.6), Inches(7.0), Inches(3.55), MUTED, 1.3)
    food_card(s, Inches(10.15), Inches(1.45), 1.9, 2.25)
    rect(s, Inches(9.6), Inches(4.6), Inches(2.75), Inches(1.3), fill=WHITE)
    text(s, Inches(9.85), Inches(4.82), Inches(2.25), Inches(0.2), "关系结论", size=13, bold=True, color=NAVY)
    text(s, Inches(9.85), Inches(5.18), Inches(2.25), Inches(0.45), "情绪集中更易互动，混合评价会削弱点赞表现。", size=10, color=INK)
    return s


def slide_baseline(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "4.4", "基准线性回归结果分析", "不放完整回归表，只保留上台最该讲的结果。")
    metric(s, Inches(0.8), Inches(1.35), "情绪密度系数", "0.0949", "正向影响", fill=TEAL_SOFT, accent=TEAL)
    metric(s, Inches(3.55), Inches(1.35), "混合评价系数", "-0.3389", "显著负向", fill=RED_SOFT, accent=RED)
    metric(s, Inches(6.3), Inches(1.35), "配图数量系数", "0.2386", "正向影响", fill=BLUE_SOFT, accent=NAVY)
    metric(s, Inches(9.05), Inches(1.35), "调整后 R²", "0.1596", "模型解释力", fill=ORANGE_SOFT, accent=ORANGE)
    rect(s, Inches(0.8), Inches(3.0), Inches(5.85), Inches(3.0), fill=WHITE)
    picture(s, FIG / "第四部分_核心变量系数图.png", Inches(1.05), Inches(3.23), w=Inches(5.35))
    rect(s, Inches(7.0), Inches(3.0), Inches(5.25), Inches(3.0), fill=WHITE)
    text(s, Inches(7.32), Inches(3.3), Inches(4.5), Inches(0.25), "一句话解释", size=15, bold=True, color=NAVY)
    bullets(s, Inches(7.32), Inches(3.75), Inches(4.55), Inches(1.8), [
        "情绪表达越集中，点赞热度整体越高。",
        "混合评价会让用户更谨慎，削弱点赞表现。",
        "配图和发布时间同样影响内容热度。"
    ], size=12)
    text(s, Inches(0.9), Inches(6.35), Inches(11.0), Inches(0.35), "关键判断：单纯“正向词更多”不一定最有效，更重要的是表达集中、态度清晰。", size=14, bold=True, color=RED)
    return s


def slide_robust_design(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "4.5", "鲁棒性检验设计", "从换标准误、换变量、换样本和换模型四个方向验证结论可靠性。")
    checks = [
        ("HC3 稳健标准误", "检验异方差影响"),
        ("替换因变量", "点赞、评论、收藏"),
        ("替换情感度量", "原始得分、正负占比"),
        ("剔除特殊样本", "极端值、凌晨、营销话术"),
        ("换模型形式", "负二项模型复测"),
    ]
    for i, (title, note) in enumerate(checks):
        x = Inches(0.95 + (i % 3) * 3.9)
        y = Inches(1.7 + (i // 3) * 2.0)
        fill = BLUE_SOFT if i % 2 == 0 else ORANGE_SOFT
        rect(s, x, y, Inches(3.25), Inches(1.35), fill=fill, line=fill)
        icon_circle(s, x + Inches(0.22), y + Inches(0.34), "✓", WHITE, TEAL)
        text(s, x + Inches(0.78), y + Inches(0.26), Inches(2.0), Inches(0.22), title, size=14, bold=True, color=NAVY)
        text(s, x + Inches(0.78), y + Inches(0.68), Inches(2.0), Inches(0.22), note, size=10, color=MUTED)
    rect(s, Inches(4.2), Inches(5.75), Inches(4.8), Inches(0.72), fill=RED_SOFT, line=RED_SOFT)
    text(s, Inches(4.35), Inches(5.94), Inches(4.5), Inches(0.22), "判断标准：核心系数方向、大小和显著性是否稳定", size=13, bold=True, color=RED, align=PP_ALIGN.CENTER)
    return s


def slide_robust_result(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "4.6", "鲁棒性检验结果", "后半部分减密度：只保留两个核心判断。")
    rect(s, Inches(0.8), Inches(1.35), Inches(5.55), Inches(4.55), fill=WHITE)
    text(s, Inches(1.15), Inches(1.65), Inches(4.7), Inches(0.28), "情绪密度：方向稳定为正", size=17, bold=True, color=TEAL)
    picture(s, OLD_FIG / "forest_density.png", Inches(1.05), Inches(2.15), w=Inches(5.05))
    text(s, Inches(1.18), Inches(5.45), Inches(4.85), Inches(0.28), "结论：有正向作用，但显著性对模型设定较敏感。", size=11, color=MUTED)
    rect(s, Inches(6.95), Inches(1.35), Inches(5.55), Inches(4.55), fill=WHITE)
    text(s, Inches(7.3), Inches(1.65), Inches(4.7), Inches(0.28), "混合评价：稳定显著为负", size=17, bold=True, color=RED)
    picture(s, OLD_FIG / "forest_mixed.png", Inches(7.2), Inches(2.15), w=Inches(5.05))
    text(s, Inches(7.32), Inches(5.45), Inches(4.85), Inches(0.28), "结论：这是本文最稳健的经验发现。", size=11, color=MUTED)
    rect(s, Inches(3.2), Inches(6.35), Inches(6.9), Inches(0.58), fill=ORANGE_SOFT, line=ORANGE_SOFT)
    text(s, Inches(3.3), Inches(6.48), Inches(6.7), Inches(0.2), "混合评价会削弱点赞热度，这一结论经得住多种检验。", size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    return s


def slide_visual_conclusion(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "4.7", "可视化结果与实证结论", "用图形把第四部分的发现收束成四个判断。")
    cards = [
        ("情绪密度", "正向促进", TEAL_SOFT, TEAL),
        ("混合评价", "负向抑制", RED_SOFT, RED),
        ("配图数量", "提升互动", BLUE_SOFT, NAVY),
        ("发布时间", "影响曝光", ORANGE_SOFT, ORANGE),
    ]
    for i, (title, note, fill, color) in enumerate(cards):
        x = Inches(0.8 + (i % 2) * 2.85)
        y = Inches(1.45 + (i // 2) * 1.35)
        rect(s, x, y, Inches(2.45), Inches(1.0), fill=fill, line=fill)
        text(s, x + Inches(0.15), y + Inches(0.18), Inches(2.1), Inches(0.18), title, size=13, bold=True, color=color)
        text(s, x + Inches(0.15), y + Inches(0.55), Inches(2.1), Inches(0.18), note, size=11, color=INK)
    rect(s, Inches(6.8), Inches(1.28), Inches(5.55), Inches(1.65), fill=WHITE)
    picture(s, FIG / "第四部分_柱状图_情绪密度与点赞.png", Inches(7.05), Inches(1.45), w=Inches(2.45))
    picture(s, FIG / "第四部分_散点图_情绪密度与点赞.png", Inches(9.55), Inches(1.45), w=Inches(2.45))
    rect(s, Inches(6.8), Inches(3.2), Inches(5.55), Inches(2.1), fill=WHITE)
    picture(s, FIG / "第四部分_折线图_发布时间与点赞.png", Inches(7.05), Inches(3.42), w=Inches(5.05))
    rect(s, Inches(0.85), Inches(5.45), Inches(11.5), Inches(1.15), fill=RED_SOFT, line=RED_SOFT)
    text(s, Inches(1.1), Inches(5.68), Inches(11.0), Inches(0.25), "最终结论", size=14, bold=True, color=RED)
    text(s, Inches(1.1), Inches(6.05), Inches(10.8), Inches(0.26), "种草笔记的点赞热度来自情感表达、评价结构与平台传播环境的共同作用。", size=15, bold=True, color=NAVY)
    return s


def slide_appendix(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "附", "答疑备用：关键变量与结果", "正式汇报可以不讲，需要时翻到这一页。")
    key = pd.read_csv(TAB / "第四部分_基准回归关键结果.csv")
    rows = key[key["变量"].isin(["情绪密度", "混合评价虚拟变量", "配图数量", "感叹号数量", "上午发布", "下午发布", "晚上发布"])][["变量", "系数", "p值", "结论"]]
    table = s.shapes.add_table(rows.shape[0] + 1, rows.shape[1], Inches(0.9), Inches(1.5), Inches(11.4), Inches(4.8)).table
    for j, col in enumerate(rows.columns):
        cell = table.cell(0, j)
        cell.text = col
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.name = "PingFang SC"
                r.font.size = Pt(10)
                r.font.bold = True
                r.font.color.rgb = WHITE
    for i in range(rows.shape[0]):
        for j in range(rows.shape[1]):
            cell = table.cell(i + 1, j)
            cell.text = str(rows.iloc[i, j])
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if i % 2 == 0 else BG
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for r in p.runs:
                    r.font.name = "PingFang SC"
                    r.font.size = Pt(9)
                    r.font.color.rgb = INK
    return s


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide_cover(prs)
    slide_framework(prs)
    slide_model(prs)
    slide_relation(prs)
    slide_baseline(prs)
    slide_robust_design(prs)
    slide_robust_result(prs)
    slide_visual_conclusion(prs)
    slide_appendix(prs)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
