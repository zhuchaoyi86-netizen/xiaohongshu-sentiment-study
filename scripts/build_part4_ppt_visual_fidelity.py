from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path("/Users/xinxinhuashe/Documents/1/python课")
OUT = ROOT / "小组共享材料" / "第四部分_实证模型与数据分析_视觉保真可编辑版.pptx"
FIG = ROOT / "第四部分材料" / "图"
OLD_FIG = ROOT / "图"


BG = RGBColor(252, 249, 243)
PAPER = RGBColor(255, 255, 252)
PAPER2 = RGBColor(255, 248, 241)
NAVY = RGBColor(22, 45, 78)
INK = RGBColor(45, 55, 68)
MUTED = RGBColor(113, 121, 132)
LINE = RGBColor(230, 221, 211)
RED = RGBColor(232, 83, 74)
RED_DARK = RGBColor(185, 55, 50)
RED_SOFT = RGBColor(255, 230, 224)
TEAL = RGBColor(45, 144, 139)
TEAL_SOFT = RGBColor(222, 244, 241)
BLUE = RGBColor(70, 119, 165)
BLUE_SOFT = RGBColor(226, 236, 249)
ORANGE = RGBColor(238, 139, 66)
ORANGE_SOFT = RGBColor(255, 239, 221)
GREEN = RGBColor(87, 137, 92)
GREEN_SOFT = RGBColor(230, 242, 225)
WHITE = RGBColor(255, 255, 255)


def set_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG


def add_text(slide, x, y, w, h, s, size=14, bold=False, color=INK, align=PP_ALIGN.LEFT, valign=MSO_VERTICAL_ANCHOR.MIDDLE):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = s
    run.font.name = "PingFang SC"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_lines(slide, x, y, w, h, items, size=10.5, color=INK, bullet=True, gap=4):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ("• " if bullet else "") + item
        p.space_after = Pt(gap)
        for run in p.runs:
            run.font.name = "PingFang SC"
            run.font.size = Pt(size)
            run.font.color.rgb = color
    return box


def shape(slide, kind, x, y, w, h, fill=PAPER, line=LINE, radius=True):
    st = kind
    shp = slide.shapes.add_shape(st, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(0.9)
    return shp


def rect(slide, x, y, w, h, fill=PAPER, line=LINE, radius=True):
    kind = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    return shape(slide, kind, x, y, w, h, fill, line, radius)


def circle(slide, x, y, d, fill, line=None):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x, y, d, d)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line:
        shp.line.color.rgb = line
    else:
        shp.line.fill.background()
    return shp


def line(slide, x1, y1, x2, y2, color=MUTED, width=1.2):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    c.line.color.rgb = color
    c.line.width = Pt(width)
    return c


def picture(slide, path, x, y, w=None, h=None):
    kwargs = {}
    if w:
        kwargs["width"] = w
    if h:
        kwargs["height"] = h
    slide.shapes.add_picture(str(path), x, y, **kwargs)


def dots(slide, x, y, rows=4, cols=6, color=LINE):
    for r in range(rows):
        for c in range(cols):
            circle(slide, x + Inches(c * 0.13), y + Inches(r * 0.13), Inches(0.025), color)


def header(slide, sec, title, subtitle=None):
    set_bg(slide)
    rect(slide, Inches(0.55), Inches(0.34), Inches(0.5), Inches(0.36), fill=RED, line=RED)
    add_text(slide, Inches(0.55), Inches(0.34), Inches(0.5), Inches(0.36), sec, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1.18), Inches(0.28), Inches(7.5), Inches(0.43), title, size=24, bold=True, color=NAVY)
    if subtitle:
        add_text(slide, Inches(1.2), Inches(0.78), Inches(8.7), Inches(0.24), subtitle, size=9.5, color=MUTED)
    dots(slide, Inches(11.45), Inches(0.38), rows=3, cols=5)
    circle(slide, Inches(12.25), Inches(0.38), Inches(0.07), RED)


def mini_food_card(slide, x, y, w=2.1, h=2.5, rotation=-3):
    base = rect(slide, x, y, Inches(w), Inches(h), fill=WHITE, line=RGBColor(238, 228, 219))
    base.rotation = rotation
    pic = rect(slide, x + Inches(0.16), y + Inches(0.16), Inches(w - 0.32), Inches(1.12), fill=ORANGE_SOFT, line=ORANGE_SOFT)
    pic.rotation = rotation
    plate = circle(slide, x + Inches(0.73), y + Inches(0.36), Inches(0.62), WHITE, RGBColor(245, 196, 161))
    plate.rotation = rotation
    for dx, dy, col in [(0.84, 0.47, RED), (1.02, 0.57, ORANGE), (1.17, 0.45, TEAL)]:
        circle(slide, x + Inches(dx), y + Inches(dy), Inches(0.11), col)
    add_text(slide, x + Inches(0.22), y + Inches(1.45), Inches(w - 0.44), Inches(0.2), "食品种草笔记", size=9, bold=True, color=NAVY)
    add_text(slide, x + Inches(0.22), y + Inches(1.77), Inches(w - 0.44), Inches(0.28), "情绪密度 · 点赞热度", size=7.5, color=MUTED)
    circle(slide, x + Inches(0.22), y + Inches(2.15), Inches(0.28), RED_SOFT)
    add_text(slide, x + Inches(0.22), y + Inches(2.15), Inches(0.28), Inches(0.28), "♥", size=10, bold=True, color=RED, align=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.58), y + Inches(2.16), Inches(0.75), Inches(0.2), "5.2k", size=8.5, bold=True, color=RED)


def metric_card(slide, x, y, label, value, note, fill=BLUE_SOFT, accent=BLUE, w=2.55):
    rect(slide, x, y, Inches(w), Inches(1.08), fill=fill, line=fill)
    add_text(slide, x + Inches(0.17), y + Inches(0.12), Inches(w - 0.34), Inches(0.18), label, size=8.5, bold=True, color=accent)
    add_text(slide, x + Inches(0.17), y + Inches(0.38), Inches(w - 0.34), Inches(0.32), value, size=20, bold=True, color=NAVY)
    add_text(slide, x + Inches(0.17), y + Inches(0.76), Inches(w - 0.34), Inches(0.18), note, size=7.7, color=MUTED)


def tiny_bar_chart(slide, x, y, w=1.7, h=0.85):
    rect(slide, x, y, Inches(w), Inches(h), fill=WHITE, line=LINE)
    vals = [0.35, 0.52, 0.7, 0.48, 0.82]
    colors = [BLUE_SOFT, TEAL_SOFT, ORANGE_SOFT, BLUE_SOFT, RED_SOFT]
    for i, v in enumerate(vals):
        bx = x + Inches(0.2 + i * 0.25)
        by = y + Inches(0.68 - v * 0.55)
        rect(slide, bx, by, Inches(0.13), Inches(v * 0.55), fill=colors[i], line=colors[i], radius=False)


def slide1(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    # left editorial block
    rect(s, Inches(0), Inches(0), Inches(4.25), Inches(7.5), fill=NAVY, line=NAVY, radius=False)
    add_text(s, Inches(0.62), Inches(0.68), Inches(1.35), Inches(0.26), "MODULE 04", size=10, bold=True, color=ORANGE_SOFT)
    add_text(s, Inches(0.62), Inches(1.28), Inches(3.05), Inches(1.65), "模块四：\n实证模型与数据分析", size=27, bold=True, color=WHITE, valign=MSO_VERTICAL_ANCHOR.TOP)
    add_text(s, Inches(0.65), Inches(3.35), Inches(2.8), Inches(0.76), "基于文本挖掘的小红书食品种草笔记情感对点赞热度的影响研究", size=11.2, color=RGBColor(238, 243, 248), valign=MSO_VERTICAL_ANCHOR.TOP)
    metric_card(s, Inches(0.65), Inches(5.35), "样本量", "2899", "有效食品种草笔记", fill=RGBColor(33, 63, 96), accent=ORANGE_SOFT, w=1.45)
    metric_card(s, Inches(2.25), Inches(5.35), "方法", "OLS", "鲁棒性检验", fill=RGBColor(33, 63, 96), accent=ORANGE_SOFT, w=1.35)
    # right visual composition
    rect(s, Inches(4.75), Inches(0.65), Inches(7.85), Inches(5.95), fill=PAPER)
    mini_food_card(s, Inches(9.72), Inches(1.03), 2.35, 2.75, rotation=4)
    for x, y, lab, fc, col in [
        (5.25, 1.65, "情绪密度", TEAL_SOFT, TEAL),
        (5.25, 2.75, "混合评价", RED_SOFT, RED),
        (7.35, 2.18, "点赞热度", ORANGE_SOFT, RED),
    ]:
        rect(s, Inches(x), Inches(y), Inches(1.45), Inches(0.5), fill=fc, line=fc)
        add_text(s, Inches(x), Inches(y), Inches(1.45), Inches(0.5), lab, size=10.5, bold=True, color=col, align=PP_ALIGN.CENTER)
    line(s, Inches(6.7), Inches(1.9), Inches(7.35), Inches(2.42), TEAL, 2)
    line(s, Inches(6.7), Inches(3.0), Inches(7.35), Inches(2.45), RED, 2)
    add_text(s, Inches(5.15), Inches(4.35), Inches(4.5), Inches(0.34), "情感表达如何影响点赞热度？", size=18, bold=True, color=NAVY)
    add_text(s, Inches(5.15), Inches(4.88), Inches(4.3), Inches(0.55), "情绪密度、混合评价、情感方向共同进入模型，并控制配图、文本长度、发布时间等变量。", size=10.5, color=INK)
    tiny_bar_chart(s, Inches(10.05), Inches(4.65), 1.75, 0.9)
    dots(s, Inches(11.6), Inches(0.9), rows=5, cols=6)


def slide2(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "4.1", "实证分析思路", "变量关系 → 模型检验 → 鲁棒性 → 可视化")
    add_text(s, Inches(0.78), Inches(1.25), Inches(8.6), Inches(0.38), "第四部分建立一条从变量到结论的证据链。", size=17, bold=True, color=NAVY)
    steps = [
        ("变量设定", "明确因变量、自变量和控制变量"),
        ("基准回归", "估计核心情感变量影响"),
        ("变量关系", "解释情绪、结构和热度联系"),
        ("鲁棒性检验", "换变量、换样本复测"),
        ("可视化展示", "用图形呈现主要关系"),
    ]
    for i, (title, note) in enumerate(steps):
        x = Inches(0.75 + i * 2.48)
        y = Inches(2.22)
        fill = RED_SOFT if i in [1, 3] else WHITE
        rect(s, x, y, Inches(1.9), Inches(2.5), fill=fill, line=LINE)
        circle(s, x + Inches(0.72), y + Inches(0.25), Inches(0.45), RED if i in [1, 3] else NAVY)
        add_text(s, x + Inches(0.72), y + Inches(0.25), Inches(0.45), Inches(0.45), f"{i+1:02d}", size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.15), y + Inches(1.0), Inches(1.6), Inches(0.22), title, size=13.5, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.22), y + Inches(1.45), Inches(1.45), Inches(0.58), note, size=8.8, color=MUTED, align=PP_ALIGN.CENTER)
        if i < 4:
            line(s, x + Inches(1.9), y + Inches(1.22), x + Inches(2.45), y + Inches(1.22), ORANGE, 1.8)
    metric_card(s, Inches(0.85), Inches(5.75), "研究样本", "2899", "2022-2024 食品种草笔记", fill=ORANGE_SOFT, accent=RED, w=2.15)
    mini_food_card(s, Inches(10.5), Inches(4.75), 1.6, 1.85, rotation=-4)


def slide3(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "4.2", "研究模型设定", "以点赞热度为因变量，以情感变量为核心自变量。")
    groups = [
        ("因变量", "log_likes_winsorized\n缩尾后的对数点赞量", RED_SOFT, RED),
        ("核心自变量", "sentiment_density\nsentiment_score_winsorized\nmixed_review_flag", BLUE_SOFT, BLUE),
        ("控制变量", "text_length_winsorized / imagenumber\nhashtag_count / time_period", GREEN_SOFT, GREEN),
    ]
    for i, (title, body, fill, color) in enumerate(groups):
        y = Inches(1.35 + i * 1.45)
        rect(s, Inches(0.75), y, Inches(4.35), Inches(1.12), fill=fill, line=fill)
        circle(s, Inches(1.03), y + Inches(0.27), Inches(0.52), WHITE)
        add_text(s, Inches(1.03), y + Inches(0.27), Inches(0.52), Inches(0.52), ["♥", "↗", "≡"][i], size=14, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(s, Inches(1.72), y + Inches(0.18), Inches(1.5), Inches(0.22), title, size=12.5, bold=True, color=color)
        add_text(s, Inches(1.72), y + Inches(0.5), Inches(2.9), Inches(0.42), body, size=9.5, color=INK, valign=MSO_VERTICAL_ANCHOR.TOP)
    rect(s, Inches(5.55), Inches(1.35), Inches(6.95), Inches(2.15), fill=WHITE)
    add_text(s, Inches(5.9), Inches(1.68), Inches(2.8), Inches(0.24), "基准回归模型（OLS）", size=15, bold=True, color=NAVY)
    formula = "log_likes_winsorized = β₀\n+ β₁ sentiment_density\n+ β₂ sentiment_score_winsorized\n+ β₃ mixed_review_flag\n+ β₄ Controls + ε"
    add_text(s, Inches(5.95), Inches(2.0), Inches(4.9), Inches(1.05), formula, size=14, bold=True, color=INK, valign=MSO_VERTICAL_ANCHOR.TOP)
    tiny_bar_chart(s, Inches(10.5), Inches(2.05), 1.45, 0.85)
    rect(s, Inches(5.55), Inches(4.08), Inches(6.95), Inches(1.45), fill=ORANGE_SOFT, line=ORANGE_SOFT)
    add_text(s, Inches(5.9), Inches(4.32), Inches(5.8), Inches(0.24), "变量处理说明", size=14, bold=True, color=RED)
    add_text(s, Inches(5.9), Inches(4.72), Inches(5.9), Inches(0.34), "点赞量先缩尾再取对数，用于降低爆款极端值对模型估计的影响。", size=10.5, color=INK)
    mini_food_card(s, Inches(10.45), Inches(5.32), 1.35, 1.48, rotation=6)


def slide4(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "4.3", "变量之间的模型关系", "点赞热度由情绪表达、评价结构和内容呈现共同影响。")
    # central relationship map
    rect(s, Inches(6.05), Inches(2.55), Inches(2.15), Inches(1.15), fill=ORANGE_SOFT, line=ORANGE_SOFT)
    add_text(s, Inches(6.05), Inches(2.78), Inches(2.15), Inches(0.28), "点赞热度", size=18, bold=True, color=RED, align=PP_ALIGN.CENTER)
    add_text(s, Inches(6.05), Inches(3.2), Inches(2.15), Inches(0.18), "log_likes", size=8.5, color=MUTED, align=PP_ALIGN.CENTER)
    nodes = [
        (0.9, 1.35, "情绪密度", "正向促进", TEAL_SOFT, TEAL, "+"),
        (0.9, 2.85, "情感得分", "结合密度解释", BLUE_SOFT, BLUE, "±"),
        (0.9, 4.35, "混合评价", "负向抑制", RED_SOFT, RED, "-"),
    ]
    for x, y, title, note, fill, color, sign in nodes:
        rect(s, Inches(x), Inches(y), Inches(2.35), Inches(0.88), fill=fill, line=fill)
        add_text(s, Inches(x + 0.18), Inches(y + 0.12), Inches(1.5), Inches(0.2), title, size=13, bold=True, color=color)
        add_text(s, Inches(x + 0.18), Inches(y + 0.45), Inches(1.5), Inches(0.18), note, size=8.7, color=MUTED)
        add_text(s, Inches(x + 1.85), Inches(y + 0.16), Inches(0.32), Inches(0.32), sign, size=16, bold=True, color=color, align=PP_ALIGN.CENTER)
    line(s, Inches(3.25), Inches(1.8), Inches(6.05), Inches(3.0), TEAL, 2)
    line(s, Inches(3.25), Inches(3.28), Inches(6.05), Inches(3.12), BLUE, 1.7)
    line(s, Inches(3.25), Inches(4.78), Inches(6.05), Inches(3.28), RED, 2)
    rect(s, Inches(4.85), Inches(5.2), Inches(4.45), Inches(0.88), fill=WHITE)
    add_text(s, Inches(5.12), Inches(5.35), Inches(3.8), Inches(0.18), "控制变量：配图数量、文本长度、发布时间、营销话术", size=10, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    line(s, Inches(7.05), Inches(5.2), Inches(7.05), Inches(3.7), MUTED, 1.3)
    mini_food_card(s, Inches(10.12), Inches(1.26), 1.85, 2.15, rotation=5)
    rect(s, Inches(9.75), Inches(4.5), Inches(2.6), Inches(1.22), fill=WHITE)
    add_text(s, Inches(9.98), Inches(4.7), Inches(2.1), Inches(0.2), "关系结论", size=12.5, bold=True, color=NAVY)
    add_text(s, Inches(9.98), Inches(5.06), Inches(2.1), Inches(0.34), "情绪集中更易互动；混合评价削弱点赞。", size=9.5, color=INK)


def slide5(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "4.4", "基准线性回归结果分析", "情绪密度正向，混合评价负向，配图数量正向。")
    metric_card(s, Inches(0.75), Inches(1.25), "情绪密度系数", "0.0949", "正向影响", fill=TEAL_SOFT, accent=TEAL, w=2.45)
    metric_card(s, Inches(3.35), Inches(1.25), "混合评价系数", "-0.3389", "显著负向", fill=RED_SOFT, accent=RED, w=2.45)
    metric_card(s, Inches(5.95), Inches(1.25), "配图数量系数", "0.2386", "正向影响", fill=BLUE_SOFT, accent=BLUE, w=2.45)
    metric_card(s, Inches(8.55), Inches(1.25), "调整后 R²", "0.1596", "模型解释力", fill=ORANGE_SOFT, accent=ORANGE, w=2.45)
    rect(s, Inches(0.75), Inches(2.95), Inches(6.25), Inches(3.05), fill=WHITE)
    picture(s, FIG / "第四部分_核心变量系数图.png", Inches(1.02), Inches(3.18), w=Inches(5.72))
    rect(s, Inches(7.35), Inches(2.95), Inches(4.8), Inches(3.05), fill=WHITE)
    add_text(s, Inches(7.68), Inches(3.25), Inches(3.9), Inches(0.25), "结果怎么讲？", size=15.5, bold=True, color=NAVY)
    add_lines(s, Inches(7.68), Inches(3.72), Inches(3.95), Inches(1.35), [
        "情绪表达越集中，点赞热度整体越高。",
        "混合评价会让用户更谨慎。",
        "配图数量和发布时间同样影响热度。"
    ], size=11.3)
    rect(s, Inches(7.68), Inches(5.22), Inches(3.85), Inches(0.45), fill=RED_SOFT, line=RED_SOFT)
    add_text(s, Inches(7.82), Inches(5.31), Inches(3.55), Inches(0.14), "不是越正向越好，而是态度越清晰越好", size=10.2, bold=True, color=RED, align=PP_ALIGN.CENTER)
    mini_food_card(s, Inches(11.08), Inches(0.86), 1.2, 1.35, rotation=6)


def slide6(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "4.5", "鲁棒性检验设计", "用多种复测方式验证结论可靠性。")
    checks = [
        ("HC3 稳健标准误", "检验异方差影响"),
        ("替换因变量", "点赞、评论、收藏"),
        ("替换情感度量", "原始得分、正负占比"),
        ("剔除特殊样本", "极端值、凌晨、营销话术"),
        ("换模型形式", "负二项模型复测"),
    ]
    for i, (title, note) in enumerate(checks):
        x = Inches(0.95 + (i % 3) * 3.85)
        y = Inches(1.55 + (i // 3) * 1.9)
        fill = BLUE_SOFT if i % 2 == 0 else ORANGE_SOFT
        rect(s, x, y, Inches(3.2), Inches(1.25), fill=fill, line=fill)
        circle(s, x + Inches(0.22), y + Inches(0.35), Inches(0.36), WHITE)
        add_text(s, x + Inches(0.22), y + Inches(0.35), Inches(0.36), Inches(0.36), "✓", size=12, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.75), y + Inches(0.22), Inches(2.15), Inches(0.22), title, size=13, bold=True, color=NAVY)
        add_text(s, x + Inches(0.75), y + Inches(0.63), Inches(2.15), Inches(0.2), note, size=9.3, color=MUTED)
    rect(s, Inches(4.15), Inches(5.62), Inches(5.05), Inches(0.72), fill=RED_SOFT, line=RED_SOFT)
    add_text(s, Inches(4.35), Inches(5.82), Inches(4.6), Inches(0.2), "判断标准：核心系数方向、大小和显著性是否稳定", size=12.5, bold=True, color=RED, align=PP_ALIGN.CENTER)
    mini_food_card(s, Inches(10.75), Inches(5.25), 1.2, 1.35, rotation=-5)


def slide7(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "4.6", "鲁棒性检验结果", "减密度展示：只保留两个核心判断。")
    rect(s, Inches(0.75), Inches(1.35), Inches(5.7), Inches(4.55), fill=WHITE)
    add_text(s, Inches(1.05), Inches(1.62), Inches(4.7), Inches(0.26), "情绪密度：方向稳定为正", size=17, bold=True, color=TEAL)
    picture(s, OLD_FIG / "forest_density.png", Inches(1.0), Inches(2.1), w=Inches(5.15))
    rect(s, Inches(1.1), Inches(5.32), Inches(4.95), Inches(0.34), fill=TEAL_SOFT, line=TEAL_SOFT)
    add_text(s, Inches(1.18), Inches(5.4), Inches(4.8), Inches(0.12), "有正向作用，但显著性对模型设定较敏感", size=9.5, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    rect(s, Inches(6.9), Inches(1.35), Inches(5.7), Inches(4.55), fill=WHITE)
    add_text(s, Inches(7.2), Inches(1.62), Inches(4.7), Inches(0.26), "混合评价：稳定显著为负", size=17, bold=True, color=RED)
    picture(s, OLD_FIG / "forest_mixed.png", Inches(7.15), Inches(2.1), w=Inches(5.15))
    rect(s, Inches(7.25), Inches(5.32), Inches(4.95), Inches(0.34), fill=RED_SOFT, line=RED_SOFT)
    add_text(s, Inches(7.33), Inches(5.4), Inches(4.8), Inches(0.12), "这是本文最稳健的经验发现", size=9.5, bold=True, color=RED, align=PP_ALIGN.CENTER)
    rect(s, Inches(3.15), Inches(6.35), Inches(7.0), Inches(0.52), fill=ORANGE_SOFT, line=ORANGE_SOFT)
    add_text(s, Inches(3.25), Inches(6.48), Inches(6.8), Inches(0.15), "混合评价会削弱点赞热度，这一结论经得住多种检验。", size=12.5, bold=True, color=NAVY, align=PP_ALIGN.CENTER)


def slide8(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "4.7", "可视化结果与实证结论", "情感表达、评价结构与平台传播环境共同作用。")
    cards = [
        ("情绪密度", "正向促进", TEAL_SOFT, TEAL),
        ("混合评价", "负向抑制", RED_SOFT, RED),
        ("配图数量", "提升互动", BLUE_SOFT, BLUE),
        ("发布时间", "影响曝光", ORANGE_SOFT, ORANGE),
    ]
    for i, (title, note, fill, color) in enumerate(cards):
        x = Inches(0.78 + (i % 2) * 2.75)
        y = Inches(1.38 + (i // 2) * 1.28)
        rect(s, x, y, Inches(2.35), Inches(0.92), fill=fill, line=fill)
        add_text(s, x + Inches(0.16), y + Inches(0.15), Inches(1.85), Inches(0.18), title, size=12.8, bold=True, color=color)
        add_text(s, x + Inches(0.16), y + Inches(0.52), Inches(1.85), Inches(0.16), note, size=10, color=INK)
    rect(s, Inches(6.45), Inches(1.22), Inches(5.9), Inches(1.75), fill=WHITE)
    picture(s, FIG / "第四部分_柱状图_情绪密度与点赞.png", Inches(6.65), Inches(1.44), w=Inches(2.7))
    picture(s, FIG / "第四部分_散点图_情绪密度与点赞.png", Inches(9.35), Inches(1.44), w=Inches(2.75))
    rect(s, Inches(6.45), Inches(3.22), Inches(5.9), Inches(2.06), fill=WHITE)
    picture(s, FIG / "第四部分_折线图_发布时间与点赞.png", Inches(6.72), Inches(3.43), w=Inches(5.35))
    rect(s, Inches(0.85), Inches(5.48), Inches(11.5), Inches(1.05), fill=RED_SOFT, line=RED_SOFT)
    add_text(s, Inches(1.12), Inches(5.68), Inches(10.8), Inches(0.24), "最终结论", size=13.5, bold=True, color=RED)
    add_text(s, Inches(1.12), Inches(6.02), Inches(10.8), Inches(0.22), "种草笔记的点赞热度来自情感表达、评价结构与平台传播环境的共同作用。", size=14.2, bold=True, color=NAVY)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    for fn in [slide1, slide2, slide3, slide4, slide5, slide6, slide7, slide8]:
        fn(prs)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
