from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path("/Users/xinxinhuashe/Documents/1/python课")
VIS_DIR = ROOT / "第四部分PPT视觉稿"
OUT = ROOT / "小组共享材料" / "第四部分_实证模型与数据分析_视觉稿高清还原版.pptx"


NAVY = RGBColor(18, 42, 74)
RED = RGBColor(232, 83, 74)
WHITE = RGBColor(255, 255, 255)
MUTED = RGBColor(100, 110, 122)


def add_text(slide, x, y, w, h, text, size=18, bold=True, color=NAVY, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "PingFang SC"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_slide_with_visual(prs, image_path, overlays):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Full-slide visual draft as the fidelity layer.
    slide.shapes.add_picture(str(image_path), 0, 0, width=prs.slide_width, height=prs.slide_height)
    # Editable text layer for the most important speaker-facing content.
    for item in overlays:
        add_text(slide, **item)
    return slide


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    overlays = {
        1: [
            dict(x=0.62, y=1.42, w=3.25, h=1.05, text="模块四：\n实证模型与数据分析", size=25, color=NAVY),
            dict(x=0.72, y=3.32, w=3.5, h=0.52, text="基于文本挖掘的小红书食品种草笔记情感对点赞热度的影响研究", size=10.5, color=RED, bold=True),
        ],
        2: [
            dict(x=0.62, y=0.82, w=4.0, h=0.52, text="4.1 实证分析思路", size=24, color=NAVY),
            dict(x=1.4, y=1.5, w=5.5, h=0.35, text="核心思路：变量关系 → 模型检验 → 鲁棒性 → 可视化", size=13, color=RED),
        ],
        3: [
            dict(x=0.62, y=0.7, w=3.8, h=0.5, text="4.2 研究模型设定", size=23, color=NAVY),
            dict(x=6.1, y=1.24, w=2.5, h=0.28, text="基准回归模型（OLS）", size=12, color=NAVY),
        ],
        4: [
            dict(x=0.62, y=0.68, w=4.8, h=0.5, text="4.3 变量之间的模型关系", size=23, color=NAVY),
            dict(x=8.4, y=1.6, w=2.8, h=0.7, text="点赞热度不是由单一情感方向决定", size=13, color=NAVY),
        ],
        5: [
            dict(x=0.62, y=0.65, w=5.2, h=0.5, text="4.4 基准线性回归结果分析", size=23, color=NAVY),
            dict(x=0.72, y=6.45, w=6.2, h=0.35, text="结论：情绪密度正向，混合评价负向，配图数量正向。", size=13, color=RED),
        ],
        6: [
            dict(x=0.62, y=0.65, w=4.3, h=0.5, text="4.5 鲁棒性检验设计", size=23, color=NAVY),
            dict(x=0.98, y=6.36, w=6.8, h=0.35, text="小结：从换标准误、换因变量、换样本、换模型形式五个角度验证可靠性。", size=11.5, color=RED),
        ],
        7: [
            dict(x=0.62, y=0.65, w=4.3, h=0.5, text="4.6 鲁棒性检验结果", size=23, color=NAVY),
            dict(x=3.2, y=6.25, w=6.8, h=0.42, text="核心结论：混合评价的负向影响最稳健。", size=14, color=RED, align=PP_ALIGN.CENTER),
        ],
        8: [
            dict(x=0.62, y=0.65, w=5.4, h=0.5, text="4.7 可视化结果与实证结论", size=23, color=NAVY),
            dict(x=0.92, y=6.45, w=9.2, h=0.35, text="本部分结论：情绪表达越集中、评价态度越清晰，越容易获得点赞。", size=13, color=RED),
        ],
    }

    for i in range(1, 9):
        img = VIS_DIR / f"第{i}页_视觉稿.png"
        add_slide_with_visual(prs, img, overlays.get(i, []))

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
